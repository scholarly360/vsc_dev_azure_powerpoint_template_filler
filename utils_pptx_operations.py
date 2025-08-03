from pptx import Presentation
import sys
import os


import json
from datetime import datetime, date
from decimal import Decimal
from typing import Any, Dict, List, Union
import inspect


"""
https://python-pptx.readthedocs.io/en/latest/user/text.html

"""

def get_font_and_text_from_slide(slide):
    """
    Analyzes fonts and text wherever applicable in a slide.


    shape -> paragraphs -> runs

    """
    text_runs = []
    if True:
        for i, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append({"run_id": i, "run_font": run.font.size, "run_text": run.text ,"run_text_len": len(run.text) })
    return text_runs

def analyze_powerpoint_file(pptx_file_path):
    """
    Analyzes a PowerPoint file to get total slides and slide layouts.
    """
    try:
        # Check if file exists
        if not os.path.exists(pptx_file_path):
            raise FileNotFoundError(f"File not found: {pptx_file_path}")
        # Load the presentation
        presentation = Presentation(pptx_file_path)
        # Analyze each slide
        slide_info = []
        
        for i, slide in enumerate(presentation.slides, 1):
            slide_layout = slide.slide_layout
            
            # Get layout information
            layout_info = {
                'slide_number': i,
                'slide_layout_type': slide_layout.name,
                'slide_layout_type_content': True if 'content' in slide_layout.name.lower() else False,
                'slide_list_text': get_font_and_text_from_slide(slide),
                'slide_original': slide
            }
            
            slide_info.append(layout_info)
        return slide_info
        
    except Exception as e:
        return {'error': str(e)}
    


