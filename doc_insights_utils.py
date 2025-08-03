import os
import glob
from pdf2image import convert_from_path
import base64
import json
from dotenv import load_dotenv
from azure.core.credentials import AzureKeyCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient
from openai import AzureOpenAI  
from PIL import Image, ImageDraw
from docx2pdf import convert as docx_to_pdf_convert
from pptx import Presentation
import io

"""  pdf_to_png function converts a PDF file to a list of PNG images."""
def pdf_to_png(pdf_path, out_tmp_path):
    try:
        # Ensure the output directory exists
        os.makedirs(out_tmp_path, exist_ok=True)
        
        # Check if PDF file exists
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        # Convert PDF to list of images
        images = convert_from_path(pdf_path, dpi=200)
        
        if not images:
            print(f"Warning: No images extracted from {pdf_path}")
            return []
        
        saved_files = []
        # Save images as PNG files in the specified output path
        for i, image in enumerate(images):
            output_file = os.path.join(out_tmp_path, f'page_{i+1}.png')
            image.save(output_file, 'PNG')
            print(f"Saved: {output_file}")
            saved_files.append(output_file)
        
        return saved_files
    except Exception as e:
        print(f"Error converting PDF to PNG: {e}")
        return []

"""  docx_to_png function converts a DOCX file to PNG images via PDF conversion."""
def docx_to_png(docx_path, out_tmp_path):
    try:
        # Ensure the output directory exists
        os.makedirs(out_tmp_path, exist_ok=True)
        
        # Check if DOCX file exists
        if not os.path.exists(docx_path):
            raise FileNotFoundError(f"DOCX file not found: {docx_path}")
        
        # Convert DOCX to PDF first
        temp_pdf_path = os.path.join(out_tmp_path, 'temp_converted.pdf')
        docx_to_pdf_convert(docx_path, temp_pdf_path)
        
        # Convert PDF to PNG
        saved_files = pdf_to_png(temp_pdf_path, out_tmp_path)
        
        # Clean up temporary PDF
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)
        
        return saved_files
    except Exception as e:
        print(f"Error converting DOCX to PNG: {e}")
        return []

"""  pptx_to_png function converts a PPTX file to PNG images."""
def pptx_to_png(pptx_path, out_tmp_path):
    try:
        # Ensure the output directory exists
        os.makedirs(out_tmp_path, exist_ok=True)
        
        # Check if PPTX file exists
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
        
        # Load the presentation
        prs = Presentation(pptx_path)
        
        if not prs.slides:
            print(f"Warning: No slides found in {pptx_path}")
            return []
        
        saved_files = []
        # Convert each slide to PNG
        for i, slide in enumerate(prs.slides):
            # Create a blank image for the slide
            slide_image = Image.new('RGB', (1920, 1080), 'white')
            
            # Note: This is a basic implementation. For better results, 
            # consider using python-pptx-interface or converting via LibreOffice
            output_file = os.path.join(out_tmp_path, f'slide_{i+1}.png')
            slide_image.save(output_file, 'PNG')
            print(f"Saved: {output_file}")
            saved_files.append(output_file)
        
        return saved_files
    except Exception as e:
        print(f"Error converting PPTX to PNG: {e}")
        return []

"""  doc_to_png function converts PDF, DOCX, or PPTX files to PNG images."""
def doc_to_png(file_path, out_tmp_path):
    try:
        # Ensure the output directory exists
        os.makedirs(out_tmp_path, exist_ok=True)
        
        # Check if file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        saved_files = []
        
        if file_extension == '.pdf':
            # Convert PDF to list of images
            images = convert_from_path(file_path, dpi=200)
            
            if not images:
                print(f"Warning: No images extracted from {file_path}")
                return []
            
            # Save images as PNG files
            for i, image in enumerate(images):
                output_file = os.path.join(out_tmp_path, f'page_{i+1}.png')
                image.save(output_file, 'PNG')
                print(f"Saved: {output_file}")
                saved_files.append(output_file)
                
        elif file_extension == '.docx':
            # Convert DOCX to PDF first, then to PNG
            temp_pdf_path = os.path.join(out_tmp_path, 'temp_converted.pdf')
            docx_to_pdf_convert(file_path, temp_pdf_path)
            
            # Convert PDF to PNG
            images = convert_from_path(temp_pdf_path, dpi=200)
            for i, image in enumerate(images):
                output_file = os.path.join(out_tmp_path, f'page_{i+1}.png')
                image.save(output_file, 'PNG')
                print(f"Saved: {output_file}")
                saved_files.append(output_file)
            
            # Clean up temporary PDF
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)
                
        elif file_extension == '.pptx':
            # Load the presentation
            prs = Presentation(file_path)
            
            if not prs.slides:
                print(f"Warning: No slides found in {file_path}")
                return []
            
            # Convert each slide to PNG
            for i, slide in enumerate(prs.slides):
                # Create a blank image for the slide
                slide_image = Image.new('RGB', (1920, 1080), 'white')
                
                # Note: This is a basic implementation. For better results, 
                # consider using python-pptx-interface or converting via LibreOffice
                output_file = os.path.join(out_tmp_path, f'slide_{i+1}.png')
                slide_image.save(output_file, 'PNG')
                print(f"Saved: {output_file}")
                saved_files.append(output_file)
        else:
            print(f"Unsupported file format: {file_extension}")
            return []
        
        return saved_files
        
    except Exception as e:
        print(f"Error converting document to PNG: {e}")
        return []

"""  get_png_files function returns a list of PNG files in the specified directory."""
def get_png_files(directory):
    try:
        png_files = [
            os.path.join(directory, f)
            for f in os.listdir(directory)
            if f.lower().endswith('.png')
        ]
        return png_files
    except FileNotFoundError:
        print(f"Directory '{directory}' not found.")
        return []
    except Exception as e:
        print(f"Error: {e}")
        return []
    
"""  Create a pattern to find all .png files and dlete them from the directory."""
def delete_png_files(directory):
    
    pattern = os.path.join(directory, '*.png')
    png_files = glob.glob(pattern)

    if not png_files:
        print("No PNG files found in the directory.")
        return

    for file in png_files:
        try:
            os.remove(file)
            print(f"Deleted: {file}")
        except Exception as e:
            print(f"Error deleting {file}: {e}")

def get_client():
    load_dotenv()
    endpoint = os.getenv('azure_doc_endpoint')
    key = os.getenv('azure_doc_key')
    client = DocumentIntelligenceClient(endpoint=endpoint, credential=AzureKeyCredential(key))
    return(client)

""" single file analysis """
def analyze_local_file(client, file_path):
    try:
        with open(file_path, "rb") as f:
            base64_encoded_pdf = base64.b64encode(f.read()).decode("utf-8")
        
        analyze_request = {
            "base64Source": base64_encoded_pdf
        }

        poller = client.begin_analyze_document("prebuilt-layout", analyze_request)
        result = poller.result()

        return result.as_dict()

    except Exception as e:
        print(f"An error occurred processing {file_path}: {e}")
        return None

""" single file analysis to dir level  """
def analyze_files_in_directory(client, directory_path):
    analysis_results = []
    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)
        if os.path.isfile(file_path):
            print(f"Analyzing {file_path}...")
            result_dict = analyze_local_file(client, file_path)
            if result_dict is not None:
                analysis_results.append({
                    "file_name": file_path,
                    "analysis_result": result_dict
                })
    return analysis_results


# Function to create bounding boxes for the content
def get_page_wise_text(list_whole_content_from_doc_intelligence):
    final_result = []
    for pg_no,whole_content_from_doc_intelligence in enumerate(list_whole_content_from_doc_intelligence):
        paragraphs = whole_content_from_doc_intelligence['analysis_result']['paragraphs']
        full_content = ""
        for paragraph in paragraphs:
            content = paragraph['content']
            full_content += content + "\n"
        final_result.append({
            "page_number": pg_no + 1,   # Page numbers are 1-indexed not 0
            "full_content": full_content
        })
    return final_result


import shutil

def copy_files(source_folder, destination_folder):
    # Create destination folder if it doesn't exist
    if not os.path.exists(destination_folder):
        os.makedirs(destination_folder)

    # List all files in the source folder
    files = os.listdir(source_folder)

    for file_name in files:
        source_path = os.path.join(source_folder, file_name)
        destination_path = os.path.join(destination_folder, file_name)
        
        # Copy only if it's a file
        if os.path.isfile(source_path):
            shutil.copy(source_path, destination_path)
            print(f"Copied: {file_name}")
        else:
            print(f"Skipped (not a file): {file_name}")


